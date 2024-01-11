import os
from pptx import Presentation

# Specify the uploads directory path
uploads_directory = "server/uploads"

# List all files in the uploads directory
pptx_files = os.listdir(uploads_directory)


def extract_text_from_pptx(pptx_file):
    # Create a Presentation object from the PowerPoint file
    prs = Presentation(pptx_file)
    all_text = []

    # Iterate through each slide and extract text from shapes
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        all_text.append(run.text)

    # Concatenate all extracted text
    return "\n".join(all_text)


if pptx_files:
    # Get the first .pptx file in the directory
    first_pptx_file = os.path.join(uploads_directory, pptx_files[0])

    # Extract text from the first .pptx file
    extracted_text = extract_text_from_pptx(first_pptx_file)

    # Remove the processed .pptx file
    os.remove(first_pptx_file)

    # Write extracted text to a new text file
    f = open("server/extracted_text.txt", "x")
    f.write(extracted_text)
    f.close()

else:
    print("No .pptx files found in the uploads directory.")
